from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "Capstone_6_Panel_Narrative_Deck.pptx"
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp"}


def find_image(candidates: Iterable[str]) -> Path | None:
    all_images = [p for p in ROOT.rglob("*") if p.suffix.lower() in IMAGE_EXTS]
    for hint in candidates:
        for img in all_images:
            if hint.lower() in img.name.lower():
                return img
    return None


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_panel_slide(
    prs: Presentation,
    title: str,
    narrative: str,
    image_path: Path | None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.6))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(28)
    title_tf.paragraphs[0].font.bold = True

    narrative_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.95), Inches(12.5), Inches(1.05))
    narrative_tf = narrative_box.text_frame
    narrative_tf.word_wrap = True
    p = narrative_tf.paragraphs[0]
    p.text = narrative
    p.font.size = Pt(15)

    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.7), Inches(2.1), width=Inches(12.0), height=Inches(4.5))
    else:
        ph = slide.shapes.add_shape(1, Inches(0.7), Inches(2.1), Inches(12.0), Inches(4.5))
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(245, 247, 250)
        ph.line.color.rgb = RGBColor(130, 130, 130)
        tf = ph.text_frame
        tf.text = "Screenshot placeholder\n(Add panel screenshot image to auto-fill next run)"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = RGBColor(90, 90, 90)


def add_summary_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.7))
    htf = header.text_frame
    htf.text = "Summary Metrics & Managerial Insights"
    htf.paragraphs[0].font.size = Pt(30)
    htf.paragraphs[0].font.bold = True

    metrics = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(5.8), Inches(2.0))
    mtf = metrics.text_frame
    mtf.text = "Summary Metrics (reference scenario)"
    mtf.paragraphs[0].font.bold = True
    mtf.paragraphs[0].font.size = Pt(18)
    for line in [
        "Annual platform cost (illustrative): ~$500K",
        "Value vs. investment (year one): 2–4× when outcomes land",
        "Targets: −10–15% time-to-degree, ~−20% excess credits, −25–30% advisor workload",
    ]:
        p = mtf.add_paragraph()
        p.text = f"- {line}"
        p.font.size = Pt(16)

    insight = slide.shapes.add_shape(1, Inches(6.7), Inches(1.1), Inches(6.0), Inches(4.8))
    insight.fill.solid()
    insight.fill.fore_color.rgb = RGBColor(234, 244, 255)
    insight.line.color.rgb = RGBColor(57, 118, 186)
    itf = insight.text_frame
    itf.word_wrap = True
    itf.text = "Conclusion (Managerial Insights)"
    itf.paragraphs[0].font.bold = True
    itf.paragraphs[0].font.size = Pt(18)
    for line in [
        "1) Treat uptime as a student-success asset: registration and pathway guidance must be dependable.",
        "2) Invest in AI + clean academic data (rules, transcripts) before scaling marketing to students.",
        "3) Train advisors and leaders together so workload shifts from triage to high-touch coaching.",
        "4) Track unit economics: ~$500K annual cost pays back when completion and efficiency gains hit target bands.",
    ]:
        p = itf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "AI Degree Advisor — Simulation Portfolio",
        "Cloud + AI academic advising: pathways, excess credits, completion — scale advising without scaling staff",
    )

    add_panel_slide(
        prs,
        "Panel 1: Budget Planning & Cost Estimation",
        "Frame ~$500K annual platform economics and multi-year ROI for AI Degree Advisor.",
        find_image(["panel1", "p1", "budget"]),
    )
    add_panel_slide(
        prs,
        "Panel 5: Platform Delivery & Automation",
        "Reliable releases for models, APIs, and integrations that students and advisors depend on.",
        find_image(["panel5", "p5", "devops"]),
    )
    add_panel_slide(
        prs,
        "Panel 6: Reliability vs Cost",
        "Balance uptime for peak registration and advising seasons against monthly cloud run rate.",
        find_image(["panel6", "p6", "uptime", "resilience"]),
    )
    add_panel_slide(
        prs,
        "Panel 8/9: Advisor & Student Adoption",
        "Advisor workflow adoption plus student diffusion — network effects accelerate student-facing rollout.",
        find_image(["panel8", "panel9", "adoption", "market"]),
    )
    add_summary_slide(prs)

    prs.save(OUTPUT)
    print(f"Created: {OUTPUT}")


if __name__ == "__main__":
    main()
